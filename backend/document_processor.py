"""Document text extraction for AISTATEweb.

Supported formats (best-effort):
  - .txt, .md: direct read (UTF-8)
  - .json: pretty-printed JSON
  - .csv: pandas -> markdown (fallback: csv module)
  - .docx: python-docx (paragraphs + tables)
  - .xlsx: openpyxl (all sheets -> markdown-ish table)
  - .pptx: python-pptx (slide texts)
  - .pdf: pdfplumber text; OCR fallback if available
  - .png/.jpg/.jpeg: pytesseract OCR (if available)

Notes:
  - OCR for PDF is optional and depends on runtime capabilities.
  - This module does not do any filesystem caching; API layer decides.
"""

from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple


SUPPORTED_EXTS = {
    ".txt",
    ".md",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".json",
    ".csv",
}


class DocumentProcessingError(RuntimeError):
    pass


@dataclass
class ExtractedDocument:
    text: str
    tables: List[Any]
    metadata: Dict[str, Any]

    def to_dict(self) -> Dict[str, Any]:
        return {
            "text": self.text,
            "tables": self.tables,
            "metadata": self.metadata,
        }


def _read_text_file(path: Path) -> ExtractedDocument:
    return ExtractedDocument(text=path.read_text(encoding="utf-8", errors="replace"), tables=[], metadata={})


def _read_json(path: Path) -> ExtractedDocument:
    try:
        obj = json.loads(path.read_text(encoding="utf-8", errors="strict"))
        text = json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception as e:
        raise DocumentProcessingError(f"Invalid JSON: {e}")
    return ExtractedDocument(text=text, tables=[], metadata={"format": "json"})


def _read_csv(path: Path) -> ExtractedDocument:
    # Prefer pandas if installed
    try:
        import pandas as pd  # type: ignore

        df = pd.read_csv(path)
        text = df.to_markdown(index=False)
        return ExtractedDocument(text=text, tables=[df.to_dict(orient="records")], metadata={"rows": int(df.shape[0]), "cols": int(df.shape[1])})
    except Exception:
        # Fallback: csv module, best-effort markdown
        with path.open("r", encoding="utf-8", errors="replace", newline="") as f:
            reader = csv.reader(f)
            rows = list(reader)
        if not rows:
            return ExtractedDocument(text="", tables=[], metadata={"rows": 0, "cols": 0})
        cols = max(len(r) for r in rows)
        # Normalize
        norm = [r + [""] * (cols - len(r)) for r in rows]
        header = norm[0]
        sep = ["---"] * cols
        md_rows = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
        for r in norm[1:]:
            md_rows.append("| " + " | ".join(r) + " |")
        return ExtractedDocument(text="\n".join(md_rows), tables=norm, metadata={"rows": len(rows) - 1, "cols": cols})


def _extract_docx(path: Path) -> ExtractedDocument:
    try:
        from docx import Document  # type: ignore
    except Exception as e:
        raise DocumentProcessingError(f"python-docx not available: {e}")

    doc = Document(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    tables_out: List[List[List[str]]] = []
    for t in doc.tables:
        table_data: List[List[str]] = []
        for row in t.rows:
            table_data.append([cell.text.strip() for cell in row.cells])
        if table_data:
            tables_out.append(table_data)

    text = "\n".join(paragraphs)
    if tables_out:
        text += "\n\n" + "\n\n".join(_table_to_md(tbl) for tbl in tables_out)

    return ExtractedDocument(text=text, tables=tables_out, metadata={"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)})


def _table_to_md(table: List[List[str]]) -> str:
    if not table:
        return ""
    cols = max(len(r) for r in table)
    norm = [r + [""] * (cols - len(r)) for r in table]
    header = norm[0]
    sep = ["---"] * cols
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(sep) + " |"]
    for r in norm[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _extract_xlsx(path: Path) -> ExtractedDocument:
    try:
        import openpyxl  # type: ignore
    except Exception as e:
        raise DocumentProcessingError(f"openpyxl not available: {e}")

    wb = openpyxl.load_workbook(str(path), data_only=True)
    sheets_text: List[str] = []
    tables: List[Any] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows: List[List[str]] = []
        max_col = ws.max_column
        max_row = ws.max_row
        # limit extremely large sheets
        max_row = min(max_row, 5000)
        max_col = min(max_col, 80)
        for r in ws.iter_rows(min_row=1, max_row=max_row, min_col=1, max_col=max_col, values_only=True):
            row = ["" if v is None else str(v) for v in r]
            # skip fully empty rows
            if any(c.strip() for c in row):
                rows.append(row)
        if rows:
            sheets_text.append(f"# Arkusz: {sheet_name}\n")
            sheets_text.append(_table_to_md(rows))
            sheets_text.append("\n")
            tables.append({"sheet": sheet_name, "rows": rows})

    return ExtractedDocument(text="\n".join(sheets_text).strip(), tables=tables, metadata={"sheets": len(wb.sheetnames)})


def _extract_pptx(path: Path) -> ExtractedDocument:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as e:
        raise DocumentProcessingError(f"python-pptx not available: {e}")

    prs = Presentation(str(path))
    parts: List[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = str(getattr(shape, "text") or "").strip()
                if txt:
                    slide_texts.append(txt)
        if slide_texts:
            parts.append(f"# Slide {i}\n" + "\n".join(slide_texts))
    return ExtractedDocument(text="\n\n".join(parts).strip(), tables=[], metadata={"slides": len(prs.slides)})


def _ocr_image(path: Path) -> ExtractedDocument:
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except Exception as e:
        raise DocumentProcessingError(f"OCR dependencies not available (Pillow+pytesseract): {e}")

    img = Image.open(str(path))
    # Best-effort languages: Polish + English
    try:
        text = pytesseract.image_to_string(img, lang="pol+eng")
    except Exception:
        text = pytesseract.image_to_string(img)
    return ExtractedDocument(text=text.strip(), tables=[], metadata={"ocr": True, "format": "image"})


def _extract_pdf(path: Path) -> ExtractedDocument:
    tables: List[Any] = []
    text_parts: List[str] = []
    pages = 0

    # 1) Try pdfplumber first
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(str(path)) as pdf:
            pages = len(pdf.pages)
            for pg in pdf.pages:
                t = (pg.extract_text() or "").strip()
                if t:
                    text_parts.append(t)
                # tables are optional and may be huge; keep only small ones
                try:
                    ptables = pg.extract_tables() or []
                    for tbl in ptables:
                        if tbl and isinstance(tbl, list) and len(tbl) <= 200:
                            tables.append(tbl)
                except Exception:
                    pass

            text = "\n\n".join(text_parts).strip()

            # 2) OCR fallback if empty
            if not text:
                try:
                    from PIL import Image  # type: ignore
                    import pytesseract  # type: ignore

                    ocr_parts: List[str] = []
                    for pg in pdf.pages:
                        try:
                            # pdfplumber's rendering backend may vary; keep it best-effort
                            im = pg.to_image(resolution=200).original  # type: ignore[attr-defined]
                            if isinstance(im, Image.Image):
                                try:
                                    ocr_parts.append(pytesseract.image_to_string(im, lang="pol+eng"))
                                except Exception:
                                    ocr_parts.append(pytesseract.image_to_string(im))
                        except Exception:
                            continue
                    text = "\n\n".join([t.strip() for t in ocr_parts if t and t.strip()]).strip()
                    if text:
                        return ExtractedDocument(text=text, tables=tables, metadata={"pages": pages, "ocr": True, "engine": "pdfplumber"})
                except Exception:
                    pass

            return ExtractedDocument(text=text, tables=tables, metadata={"pages": pages, "ocr": False, "engine": "pdfplumber"})
    except DocumentProcessingError:
        raise
    except Exception:
        # Fall through to a lighter PDF reader
        pass

    # 3) Fallback: PyPDF2 text
    try:
        import PyPDF2  # type: ignore
    except Exception as e:
        raise DocumentProcessingError(f"PDF dependencies not available (pdfplumber/PyPDF2): {e}")

    try:
        reader = PyPDF2.PdfReader(str(path))
        pages = len(reader.pages)
        for p in reader.pages:
            try:
                text_parts.append((p.extract_text() or "").strip())
            except Exception:
                pass
        text = "\n\n".join([t for t in text_parts if t]).strip()
        return ExtractedDocument(text=text, tables=[], metadata={"pages": pages, "ocr": False, "engine": "pypdf2"})
    except Exception as e:
        raise DocumentProcessingError(f"Failed to read PDF: {e}")


def extract_text(path: str | Path) -> ExtractedDocument:
    p = Path(path)
    if not p.exists() or not p.is_file():
        raise DocumentProcessingError("File does not exist")

    ext = p.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise DocumentProcessingError(f"Unsupported format: {ext}")

    if ext in (".txt", ".md"):
        return _read_text_file(p)
    if ext == ".json":
        return _read_json(p)
    if ext == ".csv":
        return _read_csv(p)
    if ext == ".docx":
        return _extract_docx(p)
    if ext == ".xlsx":
        return _extract_xlsx(p)
    if ext == ".pptx":
        return _extract_pptx(p)
    if ext in (".png", ".jpg", ".jpeg"):
        return _ocr_image(p)
    if ext == ".pdf":
        return _extract_pdf(p)

    raise DocumentProcessingError(f"Unhandled format: {ext}")
