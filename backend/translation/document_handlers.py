"""
Document handlers for various file formats
"""

import logging
import subprocess
import shutil
import tempfile
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Try to import optional dependencies
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available - DOCX support disabled")

try:
    import pdfplumber
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False
    logger.warning("pdfplumber not available - PDF support disabled")

try:
    import fitz as pymupdf  # PyMuPDF
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    logger.warning("PyMuPDF not available - PDF text overlay disabled")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available - PPTX support disabled")


class DocumentHandler:
    """Base class for document handlers"""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from document"""
        raise NotImplementedError
    
    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        **kwargs
    ):
        """Save translated text to document"""
        raise NotImplementedError


class TXTHandler(DocumentHandler):
    """Plain text file handler"""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            logger.info(f"Extracted {len(text)} characters from {file_path.name}")
            return text
        except UnicodeDecodeError:
            # Try with different encoding
            with open(file_path, 'r', encoding='latin-1') as f:
                text = f.read()
            logger.warning(f"Fallback to latin-1 encoding for {file_path.name}")
            return text
    
    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        **kwargs
    ):
        """Save translated text to TXT file"""
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(translated_text)
        logger.info(f"Saved translation to {output_path}")


class DOCXHandler(DocumentHandler):
    """Microsoft Word document handler"""

    # ------------------------------------------------------------------
    # Helpers to iterate ALL document content in reading order
    # ------------------------------------------------------------------

    @staticmethod
    def _iter_block_items(doc):
        """Yield paragraphs and tables from *doc.element.body* in document order.

        python-docx's ``doc.paragraphs`` only returns top-level paragraphs
        and completely ignores tables.  This helper walks the underlying XML
        so we get every block element (paragraph **and** table) in the order
        it appears in the document.
        """
        from docx.oxml.ns import qn  # type: ignore
        from docx.table import Table  # type: ignore
        from docx.text.paragraph import Paragraph  # type: ignore

        for child in doc.element.body:
            if child.tag == qn("w:p"):
                yield Paragraph(child, doc)
            elif child.tag == qn("w:tbl"):
                yield Table(child, doc)

    @staticmethod
    def _extract_table_text(table) -> list[str]:
        """Return a flat list of non-empty cell texts from a table."""
        texts: list[str] = []
        for row in table.rows:
            for cell in row.cells:
                ct = cell.text.strip()
                if ct:
                    texts.append(ct)
        return texts

    def extract_text(self, file_path: Path) -> str:
        """Extract text from DOCX: body paragraphs + tables + headers/footers."""
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not available")

        try:
            doc = Document(file_path)
            parts: list[str] = []

            # --- Headers (all sections) ---
            for section in doc.sections:
                for hdr in (section.header, section.first_page_header):
                    if hdr and not hdr.is_linked_to_previous:
                        for p in hdr.paragraphs:
                            if p.text.strip():
                                parts.append(p.text.strip())

            # --- Body: paragraphs + tables in document order ---
            from docx.table import Table as DocxTable  # type: ignore

            for block in self._iter_block_items(doc):
                if isinstance(block, DocxTable):
                    parts.extend(self._extract_table_text(block))
                else:
                    # Paragraph
                    if block.text.strip():
                        parts.append(block.text.strip())

            # --- Footers (all sections) ---
            for section in doc.sections:
                for ftr in (section.footer, section.first_page_footer):
                    if ftr and not ftr.is_linked_to_previous:
                        for p in ftr.paragraphs:
                            if p.text.strip():
                                parts.append(p.text.strip())

            text = "\n\n".join(parts)
            logger.info(f"Extracted {len(text)} characters from {file_path.name}")
            return text
        except Exception as e:
            logger.error(f"Failed to extract text from DOCX: {e}")
            raise
    
    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        translator=None,
        source_lang: str = None,
        target_lang: str = None,
        **kwargs
    ):
        """
        Save translated text to DOCX file
        
        If translator is provided, translates paragraph by paragraph
        to preserve formatting better
        """
        if not DOCX_AVAILABLE:
            raise RuntimeError("python-docx not available")
        
        try:
            if translator and source_lang and target_lang:
                # Paragraph-by-paragraph translation (better formatting preservation)
                doc = Document(original_path)
                
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        original_text = paragraph.text
                        translated = translator.translate_text(
                            original_text,
                            source_lang,
                            target_lang
                        )
                        
                        # Clear paragraph and add translated text
                        paragraph.clear()
                        paragraph.add_run(translated)
                
                doc.save(output_path)
            else:
                # Simple approach: create new document with translated text
                doc = Document()
                
                for paragraph_text in translated_text.split('\n\n'):
                    if paragraph_text.strip():
                        doc.add_paragraph(paragraph_text)
                
                doc.save(output_path)
            
            logger.info(f"Saved DOCX translation to {output_path}")
            
        except Exception as e:
            logger.error(f"Failed to save DOCX: {e}")
            raise


class PDFHandler(DocumentHandler):
    """PDF document handler — extract text with page markers, re-inject via PyMuPDF overlay."""

    # Unicode font for non-Latin scripts (Cyrillic, etc.)
    _FALLBACK_FONTNAME = "helv"

    # System TrueType fonts with broad Unicode coverage (tried in order)
    _UNICODE_FONT_PATHS = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
        "/usr/share/fonts/truetype/noto/NotoSans-Regular.ttf",
    ]
    _resolved_font_path: str | None = None
    _font_resolved: bool = False

    @classmethod
    def _get_unicode_fontpath(cls) -> str | None:
        """Return path to a system TTF with broad Unicode coverage, or None."""
        if cls._font_resolved:
            return cls._resolved_font_path
        cls._font_resolved = True
        for p in cls._UNICODE_FONT_PATHS:
            if Path(p).is_file():
                cls._resolved_font_path = p
                logger.info(f"PDF overlay: using Unicode font {p}")
                return p
        logger.warning("PDF overlay: no Unicode TTF found, falling back to helv (Latin only)")
        return None

    def extract_text(self, file_path: Path) -> str:
        """Extract text from PDF file, page by page with markers (like PPTX slides)."""
        if not PDF_AVAILABLE:
            raise RuntimeError("pdfplumber not available")

        try:
            parts: list[str] = []

            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    page_text = page.extract_text()
                    if page_text:
                        parts.append(f"--- Strona {i} ---\n{page_text}")
                        logger.debug(f"Extracted page {i}/{len(pdf.pages)}")

            text = "\n\n".join(parts)
            logger.info(
                f"Extracted {len(text)} chars from PDF "
                f"({len(parts)} pages, {file_path.name})"
            )
            return text

        except Exception as e:
            logger.error(f"Failed to extract text from PDF: {e}")
            raise

    # ------------------------------------------------------------------
    # PDF → PDF  text overlay (used by export-to-original endpoint)
    # ------------------------------------------------------------------

    @staticmethod
    def _normalize_color(color) -> tuple:
        """Convert PyMuPDF color (int or tuple) → (r, g, b) floats."""
        if isinstance(color, (list, tuple)):
            return tuple(color[:3])
        if isinstance(color, int):
            r = ((color >> 16) & 0xFF) / 255.0
            g = ((color >> 8) & 0xFF) / 255.0
            b = (color & 0xFF) / 255.0
            return (r, g, b)
        return (0, 0, 0)

    @staticmethod
    def _dominant_font(dict_block: dict) -> dict:
        """Return {'size': float, 'color': tuple} for the most common font in a dict block."""
        sizes: list[float] = []
        colors: list = []
        for ln in dict_block.get("lines", []):
            for sp in ln.get("spans", []):
                if sp["text"].strip():
                    sizes.append(sp["size"])
                    colors.append(sp.get("color", 0))
        if not sizes:
            return {"size": 11, "color": (0, 0, 0)}
        dominant = max(set(sizes), key=sizes.count)
        idx = sizes.index(dominant)
        return {
            "size": dominant,
            "color": PDFHandler._normalize_color(colors[idx]),
        }

    @staticmethod
    def _fit_fontsize(rect, text: str, start_size: float, min_size: float = 4.0) -> float:
        """Estimate largest font size that makes *text* fit inside *rect*.

        Uses a more accurate character-width estimation (0.55 * font_size for
        average glyph width) and accounts for line spacing (1.15x font size).
        """
        w, h = rect.width, rect.height
        if w <= 0 or h <= 0:
            return min_size
        # Try progressively smaller sizes with finer granularity
        steps = [1.0, 0.9, 0.8, 0.7, 0.6, 0.5, 0.42, 0.35, 0.28]
        for factor in steps:
            fs = max(start_size * factor, min_size)
            avg_char_w = fs * 0.55         # average character width
            line_h = fs * 1.15             # line height with spacing
            cpl = max(w / avg_char_w, 1)   # chars per line
            est_lines = 0
            for line in text.split("\n"):
                est_lines += max(1, len(line) / cpl)
            if est_lines * line_h <= h:
                return fs
        return min_size

    @staticmethod
    def _match_dict_blocks(orig_blocks: list, dict_blocks: list) -> list:
        """Match dict blocks to orig_blocks by overlapping bounding boxes.

        Instead of matching by index (which breaks if block counts differ),
        find the dict_block whose bbox overlaps most with each orig_block.
        """
        result: list[dict] = []
        for ob in orig_blocks:
            best_font = {"size": 11, "color": (0, 0, 0)}
            best_overlap = 0
            ob_rect = ob["rect"]
            for db in dict_blocks:
                db_rect = pymupdf.Rect(db["bbox"])
                overlap = ob_rect & db_rect  # intersection
                if overlap.is_empty:
                    continue
                area = overlap.width * overlap.height
                if area > best_overlap:
                    best_overlap = area
                    best_font = PDFHandler._dominant_font(db)
            result.append(best_font)
        return result

    @staticmethod
    def inject_translated_text(original_path: Path, translated_text: str) -> bytes:
        """Replace text in original PDF with *translated_text* and return PDF bytes.

        Strategy: for each page, collect all text blocks (with bounding rects),
        match translated paragraphs to blocks, then white-out originals and
        insert translated text.  Uses ``--- Strona N ---`` markers to split
        text across pages.
        """
        if not PYMUPDF_AVAILABLE:
            raise RuntimeError("PyMuPDF (fitz) not available — cannot write PDF")

        import re as _re

        doc = pymupdf.open(str(original_path))

        # --- 1. Parse translated text into pages via markers ---------------
        marker_re = _re.compile(
            r"^-{2,}\s*\S+\s+(\d+)\s*-{2,}$",
            _re.IGNORECASE | _re.MULTILINE,
        )
        pages_text: dict[int, str] = {}
        parts = marker_re.split(translated_text)
        idx = 1
        while idx + 1 < len(parts):
            page_num = int(parts[idx]) - 1
            pages_text[page_num] = parts[idx + 1].strip()
            idx += 2

        # Fallback: distribute proportionally if no markers found
        if not pages_text and translated_text.strip():
            logger.warning("PDF inject: no page markers found — distributing text across pages")
            all_lines = [ln for ln in translated_text.split("\n") if ln.strip()]
            page_char_counts: list[int] = []
            for p in doc:
                pt = p.get_text().strip()
                page_char_counts.append(len(pt) if pt else 0)
            total_chars = sum(page_char_counts) or 1
            lptr = 0
            for pi, cnt in enumerate(page_char_counts):
                if cnt == 0:
                    continue
                prop = cnt / total_chars
                n_lines = max(1, round(prop * len(all_lines)))
                if pi == len(page_char_counts) - 1:
                    pages_text[pi] = "\n".join(all_lines[lptr:])
                else:
                    pages_text[pi] = "\n".join(all_lines[lptr:lptr + n_lines])
                    lptr += n_lines

        # --- 2. Process each page ------------------------------------------
        _ufont = PDFHandler._get_unicode_fontpath()

        for page_idx, page in enumerate(doc):
            if page_idx not in pages_text:
                continue
            page_translated = pages_text[page_idx]
            if not page_translated:
                continue

            page_rect = page.rect

            # 2a. Collect ALL original text blocks before any modification
            raw_blocks = page.get_text("blocks")
            orig_blocks: list[dict] = []
            for item in raw_blocks:
                x0, y0, x1, y1, blk_text = item[0], item[1], item[2], item[3], item[4]
                btype = item[6] if len(item) > 6 else 0
                if btype != 0 or not str(blk_text).strip():
                    continue
                orig_blocks.append({
                    "rect": pymupdf.Rect(x0, y0, x1, y1),
                    "text": str(blk_text).strip(),
                })
            if not orig_blocks:
                continue

            # 2b. Get font info per block via bbox overlap matching
            dict_blocks = [
                b for b in page.get_text("dict", flags=pymupdf.TEXT_PRESERVE_WHITESPACE).get("blocks", [])
                if b.get("type", -1) == 0
            ]
            block_fonts = PDFHandler._match_dict_blocks(orig_blocks, dict_blocks)

            # 2c. Distribute translated text across blocks
            paras = [p.strip() for p in _re.split(r"\n\s*\n", page_translated) if p.strip()]
            if len(paras) == len(orig_blocks):
                chunks = paras
            else:
                all_lines = [ln for ln in page_translated.split("\n") if ln.strip()]
                total_chars = sum(len(ob["text"]) for ob in orig_blocks) or 1
                chunks: list[str] = []
                lptr = 0
                for i, ob in enumerate(orig_blocks):
                    if i == len(orig_blocks) - 1:
                        chunks.append("\n".join(all_lines[lptr:]))
                    else:
                        prop = len(ob["text"]) / total_chars
                        n = max(1, round(prop * len(all_lines)))
                        chunks.append("\n".join(all_lines[lptr:lptr + n]))
                        lptr += n

            # 2d. White-out original text — use generous padding to cover
            #     any rendering artifacts at block edges
            REDACT_PAD = 3
            has_redactions = False
            for ob in orig_blocks:
                r = ob["rect"]
                if r.is_empty or r.width < 1 or r.height < 1:
                    continue
                padded = pymupdf.Rect(
                    max(0, r.x0 - REDACT_PAD),
                    max(0, r.y0 - REDACT_PAD),
                    min(page_rect.width, r.x1 + REDACT_PAD),
                    min(page_rect.height, r.y1 + REDACT_PAD),
                )
                if not padded.is_empty:
                    page.add_redact_annot(padded, fill=(1, 1, 1))
                    has_redactions = True
            if has_redactions:
                page.apply_redactions()

            # 2e. Insert translated text block by block
            for i, ob in enumerate(orig_blocks):
                chunk = chunks[i].strip() if i < len(chunks) else ""
                if not chunk:
                    continue

                rect = ob["rect"]
                fi = block_fonts[i] if i < len(block_fonts) else {"size": 11, "color": (0, 0, 0)}

                # Validate rect — skip degenerate ones
                if rect.is_empty or rect.width < 5 or rect.height < 5:
                    logger.debug(f"PDF page {page_idx+1} block {i}: skipping degenerate rect {rect}")
                    continue

                # Start with the original font size
                orig_fs = max(fi["size"], 5.0)
                fontsize = PDFHandler._fit_fontsize(rect, chunk, orig_fs)

                # If text overflows at current rect, try expanding the rect
                # within reasonable page bounds before shrinking font further
                avg_char_w = fontsize * 0.55
                line_h = fontsize * 1.15
                cpl = max(rect.width / avg_char_w, 1)
                est_lines = sum(max(1, len(ln) / cpl) for ln in chunk.split("\n"))

                if est_lines * line_h > rect.height:
                    max_x1 = min(page_rect.width - 20, max(rect.x1, page_rect.width * 0.85))
                    # Find next block's top edge to limit vertical expansion
                    next_y0 = page_rect.height - 10
                    for j, ob2 in enumerate(orig_blocks):
                        if j > i and ob2["rect"].y0 > rect.y0 + 5:
                            next_y0 = ob2["rect"].y0 - 2
                            break
                    max_y1 = min(next_y0, rect.y1 + rect.height * 0.5)
                    max_y1 = max(max_y1, rect.y1)  # never shrink

                    expanded = pymupdf.Rect(rect.x0, rect.y0, max_x1, max_y1)
                    if not expanded.is_empty and expanded.width >= 5 and expanded.height >= 5:
                        fontsize = PDFHandler._fit_fontsize(expanded, chunk, orig_fs)
                        rect = expanded

                # Final safety: ensure rect is valid before insert
                if rect.is_empty or rect.width < 2 or rect.height < 2:
                    logger.warning(f"PDF page {page_idx+1} block {i}: rect too small after expansion, skipping")
                    continue

                tb_args: dict = dict(
                    rect=rect,
                    buffer=chunk,
                    fontsize=fontsize,
                    color=fi["color"],
                    align=pymupdf.TEXT_ALIGN_LEFT,
                )
                if _ufont:
                    tb_args["fontfile"] = _ufont
                    tb_args["fontname"] = "ujfont"

                try:
                    rc = page.insert_textbox(**tb_args)
                    if rc < 0:
                        # Text didn't fully fit — try with smaller font
                        smaller = max(4.0, fontsize * 0.7)
                        tb_args["fontsize"] = smaller
                        page.insert_textbox(**tb_args)
                        logger.debug(
                            f"PDF page {page_idx + 1} block {i}: text overflow, "
                            f"reduced font {fontsize:.1f}→{smaller:.1f}"
                        )
                except Exception as e:
                    logger.warning(
                        f"PDF page {page_idx+1} block {i}: insert_textbox failed "
                        f"(rect={rect}, fontsize={fontsize:.1f}): {e}"
                    )

        pdf_bytes = doc.tobytes()
        doc.close()
        return pdf_bytes

    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        **kwargs,
    ):
        """Save translated PDF — overlay text on original preserving layout."""
        if PYMUPDF_AVAILABLE:
            try:
                pdf_bytes = self.inject_translated_text(original_path, translated_text)
                output_path = output_path.with_suffix(".pdf")
                output_path.write_bytes(pdf_bytes)
                logger.info(f"Saved PDF translation to {output_path}")
                return
            except Exception as e:
                logger.error(f"PyMuPDF overlay failed, falling back to TXT: {e}")

        # Fallback: save as TXT
        txt_output = output_path.with_suffix(".txt")
        with open(txt_output, "w", encoding="utf-8") as f:
            f.write(translated_text)
        logger.info(f"Saved PDF translation as TXT: {txt_output}")
        logger.warning("PDF → PDF translation requires PyMuPDF library")


class SRTHandler(DocumentHandler):
    """Subtitle file handler (preserves timestamps)"""
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from SRT file (returns full content with timestamps)"""
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        logger.info(f"Extracted {len(text)} characters from SRT")
        return text
    
    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        translator=None,
        source_lang: str = None,
        target_lang: str = None,
        **kwargs
    ):
        """
        Save translated SRT file (preserving structure)
        
        If translator is provided, translates subtitle-by-subtitle
        """
        import re
        
        try:
            with open(original_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Parse SRT blocks
            blocks = re.split(r'\n\n+', content.strip())
            translated_blocks = []
            
            for block in blocks:
                lines = block.split('\n')
                
                if len(lines) >= 3:
                    # lines[0] = number
                    # lines[1] = timestamp
                    # lines[2:] = subtitle text
                    
                    number = lines[0]
                    timestamp = lines[1]
                    subtitle_text = ' '.join(lines[2:])
                    
                    # Translate subtitle
                    if translator and source_lang and target_lang:
                        translated = translator.translate_text(
                            subtitle_text,
                            source_lang,
                            target_lang
                        )
                    else:
                        translated = subtitle_text
                    
                    translated_block = f"{number}\n{timestamp}\n{translated}"
                    translated_blocks.append(translated_block)
            
            # Save
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n\n'.join(translated_blocks))
            
            logger.info(f"Saved SRT translation to {output_path}")
            
        except Exception as e:
            logger.error(f"Failed to save SRT: {e}")
            raise


class PPTXHandler(DocumentHandler):
    """Microsoft PowerPoint handler — extracts text from all slides."""

    def extract_text(self, file_path: Path) -> str:
        """Extract text from PPTX file, slide by slide."""
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available")

        try:
            prs = Presentation(file_path)
            parts: list[str] = []

            for idx, slide in enumerate(prs.slides, 1):
                slide_texts: list[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            txt = para.text.strip()
                            if txt:
                                slide_texts.append(txt)
                    # Also extract text from tables
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                txt = cell.text.strip()
                                if txt:
                                    slide_texts.append(txt)
                if slide_texts:
                    parts.append(f"--- Slajd {idx} ---\n" + "\n".join(slide_texts))

            text = "\n\n".join(parts)
            logger.info(
                f"Extracted {len(text)} chars from PPTX "
                f"({len(prs.slides)} slides, {file_path.name})"
            )
            return text

        except Exception as e:
            logger.error(f"Failed to extract text from PPTX: {e}")
            raise

    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        translator=None,
        source_lang: str = None,
        target_lang: str = None,
        **kwargs,
    ):
        """
        Save translated PPTX — clones the original presentation and replaces
        text in every shape while preserving slide layout, images, formatting.
        """
        if not PPTX_AVAILABLE:
            raise RuntimeError("python-pptx not available")

        try:
            prs = Presentation(original_path)

            if translator and source_lang and target_lang:
                # Translate shape-by-shape preserving formatting
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                full = para.text.strip()
                                if not full:
                                    continue
                                translated = translator.translate_text(
                                    full, source_lang, target_lang
                                )
                                # Overwrite first run, clear the rest
                                if para.runs:
                                    para.runs[0].text = translated
                                    for run in para.runs[1:]:
                                        run.text = ""
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    txt = cell.text.strip()
                                    if txt:
                                        translated = translator.translate_text(
                                            txt, source_lang, target_lang
                                        )
                                        cell.text = translated
                prs.save(output_path)
            else:
                # Fallback: save plain text as DOCX
                if DOCX_AVAILABLE:
                    from docx import Document as DocxDoc

                    doc = DocxDoc()
                    for para_text in translated_text.split("\n\n"):
                        if para_text.strip():
                            doc.add_paragraph(para_text)
                    doc.save(output_path.with_suffix(".docx"))
                else:
                    with open(output_path.with_suffix(".txt"), "w", encoding="utf-8") as f:
                        f.write(translated_text)

            logger.info(f"Saved PPTX translation to {output_path}")

        except Exception as e:
            logger.error(f"Failed to save PPTX translation: {e}")
            raise


class DOCHandler(DocumentHandler):
    """Microsoft Word 97-2003 (.doc) handler — converts to .docx via LibreOffice, then delegates."""

    _LIBREOFFICE_BIN: str | None = None
    _lo_checked: bool = False

    @classmethod
    def _find_libreoffice(cls) -> str | None:
        """Find LibreOffice binary on the system."""
        if cls._lo_checked:
            return cls._LIBREOFFICE_BIN
        cls._lo_checked = True
        for name in ("libreoffice", "soffice"):
            path = shutil.which(name)
            if path:
                cls._LIBREOFFICE_BIN = path
                logger.info(f"DOCHandler: using LibreOffice at {path}")
                return path
        logger.warning("DOCHandler: LibreOffice not found — .doc support unavailable")
        return None

    @staticmethod
    def convert_doc_to_docx(doc_path: Path) -> Path:
        """Convert .doc file to .docx using LibreOffice.  Returns path to the new .docx."""
        lo = DOCHandler._find_libreoffice()
        if not lo:
            raise RuntimeError(
                "LibreOffice is required to process .doc files. "
                "Please install LibreOffice or convert the file to .docx manually."
            )

        tmp_dir = tempfile.mkdtemp(prefix="aistate_doc_")
        try:
            result = subprocess.run(
                [
                    lo,
                    "--headless",
                    "--convert-to", "docx",
                    "--outdir", tmp_dir,
                    str(doc_path),
                ],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if result.returncode != 0:
                raise RuntimeError(f"LibreOffice conversion failed: {result.stderr[:500]}")

            # Find the converted file
            docx_files = list(Path(tmp_dir).glob("*.docx"))
            if not docx_files:
                raise RuntimeError("LibreOffice did not produce a .docx file")

            # Move the converted file next to the original
            out_path = doc_path.with_suffix(".docx")
            # Avoid overwriting if a .docx already exists at that path
            if out_path.exists():
                out_path = doc_path.parent / f"{doc_path.stem}_converted.docx"
            shutil.move(str(docx_files[0]), str(out_path))
            logger.info(f"Converted {doc_path.name} → {out_path.name}")
            return out_path
        except subprocess.TimeoutExpired:
            raise RuntimeError("LibreOffice conversion timed out (120s)")
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def extract_text(self, file_path: Path) -> str:
        """Extract text from .doc by converting to .docx first."""
        docx_path = self.convert_doc_to_docx(file_path)
        try:
            docx_handler = DOCXHandler()
            return docx_handler.extract_text(docx_path)
        finally:
            # Keep the converted .docx — it will be used for save_translated / export-to-original
            pass

    def save_translated(
        self,
        original_path: Path,
        translated_text: str,
        output_path: Path,
        **kwargs,
    ):
        """Save translated .doc — converts to .docx, injects text, saves as .docx."""
        docx_path = original_path.with_suffix(".docx")
        if not docx_path.exists():
            docx_path = self.convert_doc_to_docx(original_path)
        docx_handler = DOCXHandler()
        out = output_path.with_suffix(".docx")
        docx_handler.save_translated(docx_path, translated_text, out, **kwargs)


# Factory function
def get_handler(file_extension: str) -> Optional[DocumentHandler]:
    """
    Get appropriate handler for file type
    
    Args:
        file_extension: File extension (e.g., '.txt', '.docx')
        
    Returns:
        DocumentHandler instance or None if unsupported
    """
    handlers = {
        '.txt': TXTHandler,
        '.doc': DOCHandler,
        '.docx': DOCXHandler,
        '.pdf': PDFHandler,
        '.srt': SRTHandler,
        '.pptx': PPTXHandler,
    }

    ext = file_extension.lower()
    handler_class = handlers.get(ext)

    if handler_class:
        # Check if dependencies are available
        if ext == '.docx' and not DOCX_AVAILABLE:
            logger.error("DOCX handler requested but python-docx not available")
            return None
        if ext == '.pdf' and not PDF_AVAILABLE:
            logger.error("PDF handler requested but pdfplumber not available")
            return None
        if ext == '.pptx' and not PPTX_AVAILABLE:
            logger.error("PPTX handler requested but python-pptx not available")
            return None
        
        return handler_class()
    
    logger.warning(f"No handler available for extension: {ext}")
    return None


# Convenience function
def extract_text_from_file(file_path: Path) -> str:
    """
    Extract text from any supported file format
    
    Args:
        file_path: Path to file
        
    Returns:
        Extracted text
        
    Raises:
        ValueError: If file type not supported
        RuntimeError: If extraction fails
    """
    handler = get_handler(file_path.suffix)
    
    if not handler:
        raise ValueError(f"Unsupported file type: {file_path.suffix}")
    
    return handler.extract_text(file_path)
